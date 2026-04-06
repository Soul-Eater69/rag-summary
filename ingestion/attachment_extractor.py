"""
Binary attachment extractor (V6).

Accepts raw file bytes and extracts normalized text + metadata, delegating
to the best available parser for each file type.

Supported types (all with graceful fallback):
  .pdf   — pdfplumber (preferred) → PyPDF2 → text stub
  .xlsx  — openpyxl → raw bytes stub
  .xls   — xlrd → raw bytes stub
  .docx  — python-docx → raw bytes stub
  .pptx  — python-pptx → raw bytes stub
  .csv   — stdlib csv
  .txt   — direct decode

Output dict (AttachmentExtractor.extract):
  filename:          str
  file_type:         "pdf"|"xlsx"|"xls"|"docx"|"pptx"|"csv"|"txt"|"unknown"
  text:              str   — extracted, normalized text (may be partial)
  page_count:        int   — pages (PDF) or sheets (XLSX)
  sheet_names:       List[str]  — XLSX/XLS sheet names
  extraction_quality: "high"|"medium"|"low"|"failed"
  warnings:          List[str]

Callers pass _attachment_contents as:
  [{"filename": "budget.xlsx", "content": <bytes>}, ...]

The pipeline node (node_parse_attachments) calls this before AttachmentParser.
"""

from __future__ import annotations

import csv
import io
import logging
import os
from typing import Any, Dict, List

logger = logging.getLogger(__name__)

# Maximum characters to extract per attachment (prevent state bloat)
_MAX_CHARS = 20_000


class AttachmentExtractor:
    """
    File-type-aware text extractor for binary attachment files.

    All methods degrade gracefully if the required library is missing,
    returning extraction_quality="low" or "failed" with a warning.
    """

    def extract(self, filename: str, content: bytes) -> Dict[str, Any]:
        """
        Extract text and metadata from binary file content.

        Args:
            filename: Original filename (used for type detection).
            content:  Raw file bytes.

        Returns:
            dict with keys: filename, file_type, text, page_count,
            sheet_names, extraction_quality, warnings.
        """
        ext = os.path.splitext(filename.lower())[1].lstrip(".")
        dispatch = {
            "pdf": self._extract_pdf,
            "xlsx": self._extract_xlsx,
            "xls": self._extract_xls,
            "docx": self._extract_docx,
            "pptx": self._extract_pptx,
            "csv": self._extract_csv,
            "txt": self._extract_txt,
            "text": self._extract_txt,
        }
        extractor = dispatch.get(ext, self._extract_unknown)
        result = extractor(filename, content)
        result["filename"] = filename
        result["file_type"] = ext if ext in dispatch else "unknown"
        # Truncate to prevent oversized state
        if len(result.get("text", "")) > _MAX_CHARS:
            result["text"] = result["text"][:_MAX_CHARS]
            result.setdefault("warnings", []).append(
                f"Text truncated to {_MAX_CHARS} chars"
            )
        return result

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def _extract_pdf(self, filename: str, content: bytes) -> Dict[str, Any]:
        # Try pdfplumber first (best quality — handles tables, multi-column)
        try:
            import pdfplumber  # type: ignore
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                pages = []
                for page in pdf.pages:
                    txt = page.extract_text() or ""
                    if txt.strip():
                        pages.append(txt)
                text = "\n\n".join(pages)
                return {
                    "text": text,
                    "page_count": len(pdf.pages),
                    "sheet_names": [],
                    "extraction_quality": "high" if text.strip() else "low",
                    "warnings": [],
                }
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("[AttachmentExtractor] pdfplumber failed for %s: %s", filename, exc)

        # Fallback: PyPDF2
        try:
            import PyPDF2  # type: ignore
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            pages = []
            for page in reader.pages:
                txt = page.extract_text() or ""
                if txt.strip():
                    pages.append(txt)
            text = "\n\n".join(pages)
            return {
                "text": text,
                "page_count": len(reader.pages),
                "sheet_names": [],
                "extraction_quality": "medium" if text.strip() else "low",
                "warnings": ["pdfplumber unavailable; used PyPDF2"],
            }
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("[AttachmentExtractor] PyPDF2 failed for %s: %s", filename, exc)

        return {
            "text": "",
            "page_count": 0,
            "sheet_names": [],
            "extraction_quality": "failed",
            "warnings": ["PDF parsing requires pdfplumber or PyPDF2"],
        }

    # ------------------------------------------------------------------
    # XLSX
    # ------------------------------------------------------------------

    def _extract_xlsx(self, filename: str, content: bytes) -> Dict[str, Any]:
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            text_parts = []
            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    row_text = "\t".join(cells).strip()
                    if row_text:
                        text_parts.append(row_text)
            text = "\n".join(text_parts)
            return {
                "text": text,
                "page_count": len(sheet_names),
                "sheet_names": sheet_names,
                "extraction_quality": "high" if text.strip() else "low",
                "warnings": [],
            }
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("[AttachmentExtractor] openpyxl failed for %s: %s", filename, exc)

        return {
            "text": "",
            "page_count": 0,
            "sheet_names": [],
            "extraction_quality": "failed",
            "warnings": ["XLSX parsing requires openpyxl"],
        }

    # ------------------------------------------------------------------
    # XLS (legacy Excel)
    # ------------------------------------------------------------------

    def _extract_xls(self, filename: str, content: bytes) -> Dict[str, Any]:
        try:
            import xlrd  # type: ignore
            wb = xlrd.open_workbook(file_contents=content)
            text_parts = []
            for sheet_name in wb.sheet_names():
                ws = wb.sheet_by_name(sheet_name)
                text_parts.append(f"[Sheet: {sheet_name}]")
                for row_idx in range(ws.nrows):
                    cells = [str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)]
                    row_text = "\t".join(cells).strip()
                    if row_text:
                        text_parts.append(row_text)
            text = "\n".join(text_parts)
            return {
                "text": text,
                "page_count": wb.nsheets,
                "sheet_names": wb.sheet_names(),
                "extraction_quality": "high" if text.strip() else "low",
                "warnings": [],
            }
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("[AttachmentExtractor] xlrd failed for %s: %s", filename, exc)

        return {
            "text": "",
            "page_count": 0,
            "sheet_names": [],
            "extraction_quality": "failed",
            "warnings": ["XLS parsing requires xlrd"],
        }

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _extract_docx(self, filename: str, content: bytes) -> Dict[str, Any]:
        try:
            import docx  # type: ignore  (python-docx)
            doc = docx.Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract table cells
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        paragraphs.append("\t".join(cells))
            text = "\n".join(paragraphs)
            return {
                "text": text,
                "page_count": 1,  # DOCX doesn't expose page count directly
                "sheet_names": [],
                "extraction_quality": "high" if text.strip() else "low",
                "warnings": [],
            }
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("[AttachmentExtractor] python-docx failed for %s: %s", filename, exc)

        return {
            "text": "",
            "page_count": 0,
            "sheet_names": [],
            "extraction_quality": "failed",
            "warnings": ["DOCX parsing requires python-docx"],
        }

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _extract_pptx(self, filename: str, content: bytes) -> Dict[str, Any]:
        try:
            from pptx import Presentation  # type: ignore  (python-pptx)
            prs = Presentation(io.BytesIO(content))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                if parts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(parts))
            text = "\n\n".join(slides)
            return {
                "text": text,
                "page_count": len(prs.slides),
                "sheet_names": [],
                "extraction_quality": "high" if text.strip() else "low",
                "warnings": [],
            }
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("[AttachmentExtractor] python-pptx failed for %s: %s", filename, exc)

        return {
            "text": "",
            "page_count": 0,
            "sheet_names": [],
            "extraction_quality": "failed",
            "warnings": ["PPTX parsing requires python-pptx"],
        }

    # ------------------------------------------------------------------
    # CSV
    # ------------------------------------------------------------------

    def _extract_csv(self, filename: str, content: bytes) -> Dict[str, Any]:
        try:
            text_io = io.StringIO(content.decode("utf-8", errors="replace"))
            reader = csv.reader(text_io)
            rows = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
            text = "\n".join(rows)
            return {
                "text": text,
                "page_count": 1,
                "sheet_names": [],
                "extraction_quality": "high" if text.strip() else "low",
                "warnings": [],
            }
        except Exception as exc:
            logger.warning("[AttachmentExtractor] CSV failed for %s: %s", filename, exc)
            return {
                "text": "",
                "page_count": 0,
                "sheet_names": [],
                "extraction_quality": "failed",
                "warnings": [str(exc)],
            }

    # ------------------------------------------------------------------
    # Plain text
    # ------------------------------------------------------------------

    def _extract_txt(self, filename: str, content: bytes) -> Dict[str, Any]:
        for encoding in ("utf-8", "latin-1", "cp1252"):
            try:
                text = content.decode(encoding)
                return {
                    "text": text,
                    "page_count": 1,
                    "sheet_names": [],
                    "extraction_quality": "high",
                    "warnings": [],
                }
            except UnicodeDecodeError:
                continue
        return {
            "text": content.decode("utf-8", errors="replace"),
            "page_count": 1,
            "sheet_names": [],
            "extraction_quality": "medium",
            "warnings": ["Encoding auto-detected; some characters may be garbled"],
        }

    # ------------------------------------------------------------------
    # Unknown type fallback
    # ------------------------------------------------------------------

    def _extract_unknown(self, filename: str, content: bytes) -> Dict[str, Any]:
        # Try UTF-8 text decode as last resort
        try:
            text = content.decode("utf-8")
            if text.isprintable() or "\n" in text:
                return {
                    "text": text,
                    "page_count": 1,
                    "sheet_names": [],
                    "extraction_quality": "low",
                    "warnings": ["Unknown file type; attempted UTF-8 text decode"],
                }
        except Exception:
            pass
        return {
            "text": "",
            "page_count": 0,
            "sheet_names": [],
            "extraction_quality": "failed",
            "warnings": [f"Unsupported file type; cannot extract text"],
        }
